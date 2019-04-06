from pptx import Presentation
from pptx.util import Inches
import csv

inputfile="split.csv"

prs = Presentation('base.pptx')

source = csv.reader(open(inputfile))

for row in source:
    zh = row[0]
    en = row[1]
    
    blank_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(blank_slide_layout)
    print("生成 " + zh)
    zh_holder = slide.placeholders[0]
    zh_holder.text = zh
    en_holder = slide.placeholders[1]
    en_holder.text = en

prs.save("output.pptx")